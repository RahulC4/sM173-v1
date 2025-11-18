# ============================================================
# generate_ppt.py  (FINAL – with Azure fixes)
# ============================================================

import os
import tempfile
import uuid
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# Ensure directory exists for design JSONs
ensure_dir("design_jsons")


CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ---------------------------------------------------------
# AUTO–DETECTION: slide count + theme
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    # detect "5 slides"
    match = re.search(r"(\d+)\s+slides?", prompt.lower())
    if match:
        num_slides = int(match.group(1))

    # detect themes
    theme_words = [
        "corporate", "modern", "minimal", "professional",
        "dark", "light", "colorful", "gradient", "flat"
    ]
    for t in theme_words:
        if t in prompt.lower():
            theme = t.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# GPT PLAN GENERATOR
# ---------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):

    sys_prompt = (
        "You are an expert presentation designer.\n"
        "Create a structured slide plan.\n"
        "Return strictly JSON: a list of slides:\n"
        "[ {\"title\": str, \"bullets\": [str], "
        "\"visual_required\": bool, \"visual_prompt\": str } ]\n\n"
        "Use the design JSONs for style cues and reference snippets for content.\n"
        "If both match well, enable visuals. Otherwise, decide independently.\n"
        "If slide count or theme is detected from user text, follow it.\n"
    )

    if theme:
        sys_prompt += f"\nTheme preference: {theme}\n"

    # include truncated design + reference text
    sys_prompt += f"\nDesign Context (truncated): {json.dumps(design_context)[:4000]}"
    sys_prompt += f"\nReference Snippets: {json.dumps(references_text)[:2000]}\n"

    user_prompt = f"Create a PowerPoint plan for: {prompt}. Style: {style}."
    if num_slides:
        user_prompt += f" Produce exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_completion_tokens=1500,     # FIXED FOR AZURE
            temperature=1
        )
        raw = resp.choices[0].message.content
        plan = safe_json_load(raw)

        if not plan:
            logger.warning("Invalid plan JSON received from GPT.")
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]

        return plan

    except Exception as e:
        logger.exception(f"Plan generation failed: {e}")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]


# ---------------------------------------------------------
# VISUAL GENERATOR (Image model)
# ---------------------------------------------------------
def generate_visual_image(prompt: str):
    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024"
        )
        url = resp.data[0].url
        img_bytes = requests.get(url, timeout=20).content

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        return tmp.name

    except Exception as e:
        logger.warning(f"Image generation failed for '{prompt}': {e}")
        return None


# ---------------------------------------------------------
# PPT BUILDER
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()

    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # title
        try:
            slide.shapes.title.text = s.get("title", "")
        except:
            pass

        # bullets
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
        except:
            pass

        # image
        if s.get("image_path"):
            try:
                slide.shapes.add_picture(
                    s["image_path"],
                    Inches(0.5), Inches(3.0),
                    width=Inches(8)
                )
            except Exception as e:
                logger.debug(f"Failed to add image: {e}")

    out = os.path.join(
        tempfile.gettempdir(),
        f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx"
    )
    prs.save(out)
    return out


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None):

    detected_slides, detected_theme = parse_user_intent(prompt)
    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    refs = semantic_search(prompt, top_k=10, tags=tag_filters) or []

    design_context = []
    reference_text = []

    for r in refs:
        ppt_name = r.get("ppt_name")
        text_snip = (r.get("text") or "")[:500]
        if text_snip:
            reference_text.append(text_snip)

        json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets.")

    # GPT plan
    plan = call_llm_plan(
        prompt,
        style,
        design_context,
        reference_text,
        num_slides=requested_num_slides,
        theme=theme
    )

    # Generate images if required
    slides = []
    for sp in plan:
        slide = {
            "title": sp.get("title", "Untitled"),
            "bullets": sp.get("bullets", []),
            "image_path": None
        }

        if sp.get("visual_required"):
            img_prompt = sp.get("visual_prompt", f"Professional visual for {slide['title']}")
            slide["image_path"] = generate_visual_image(img_prompt)

        slides.append(slide)

    # Build PPT
    out_path = build_ppt(slides)

    # Upload to Azure Blob
    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    # Log metadata
    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "style": style,
        "theme": theme,
        "requested_num_slides": requested_num_slides,
        "references_used": len(reference_text),
        "design_jsons_used": len(design_context),
        "slides_generated": len(slides),
        "ppt_file": fname
    }

    upload_json_to_blob(json.dumps(log, indent=2).encode("utf-8"),
                        f"logs/{fname}.json")

    return out_path, log