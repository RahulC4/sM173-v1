import os
import uuid
import tempfile
from pptx import Presentation
from azure.storage.blob import BlobServiceClient
from openai import AzureOpenAI
from chromadb import PersistentClient
from utils import get_env, logger, now_ts, get_embedding_dim
from design_extractor import extract_design_json  # if you have it

# === TEXT client (embeddings) ===
text_client = AzureOpenAI(
    azure_endpoint=get_env("OPENAI_API_BASE", required=True),
    api_key=get_env("OPENAI_API_KEY", required=True),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

EMBEDDING_MODEL = get_env("EMBEDDING_MODEL", "text-embedding-3-large")
EMBEDDING_DIM = get_embedding_dim(EMBEDDING_MODEL)
CHROMA_PERSIST_DIR = get_env("CHROMA_PERSIST_DIR", "./chroma_db")

# === Blob client ===
blob_client = BlobServiceClient.from_connection_string(get_env("AZURE_BLOB_CONN"))
container_client = blob_client.get_container_client(get_env("AZURE_BLOB_CONTAINER"))

# === Chroma ===
chroma_client = PersistentClient(path=CHROMA_PERSIST_DIR)
try:
    collection = chroma_client.get_collection("ppt_slides")
except:
    collection = chroma_client.create_collection("ppt_slides")


def azure_embed_func(texts):
    try:
        resp = text_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=texts
        )
        return [d.embedding for d in resp.data]
    except Exception as e:
        logger.exception(f"Embedding failed: {e}")
        return []


def extract_slides(local_path):
    prs = Presentation(local_path)
    slides = []
    for i, s in enumerate(prs.slides):
        texts = []
        for shape in s.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"index": i, "text": "\n".join(texts)})
    return slides


def process_blob(blob_name):
    logger.info(f"Processing blob: {blob_name}")
    tmp = os.path.join(tempfile.gettempdir(), blob_name.replace("/", "_"))
    with open(tmp, "wb") as fp:
        container_client.download_blob(blob_name).readinto(fp)

    slides = extract_slides(tmp)
    if not slides:
        return

    docs, metas, ids, alltexts = [], [], [], []
    ppt_base = os.path.splitext(os.path.basename(blob_name))[0]

    for s in slides:
        sid = f"{ppt_base}_Slide_{s['index']}"
        text = s.get("text", "")
        meta = {
            "ppt_name": blob_name,
            "slide_id": sid,
            "slide_index": s["index"],
            "title": text.split("\n", 1)[0] if text else ""
        }

        ids.append(str(uuid.uuid4()))
        docs.append(text)
        metas.append(meta)
        alltexts.append(text)

    embeddings = azure_embed_func(alltexts)
    if not embeddings:
        return

    collection.add(
        documents=docs,
        embeddings=embeddings,
        metadatas=metas,
        ids=ids
    )

    logger.info(f"Indexed {len(slides)} slides from {blob_name}")