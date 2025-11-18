import streamlit as st 
from generate_ppt import generate_presentation
from azure_blob_utils import list_generated_presentations
from utils import logger
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import os

st.set_page_config(page_title="AI PPT Generator (Chroma)", layout="wide")
st.title("🧠 AI PowerPoint Generator (Chroma + Azure OpenAI)")

col1, col2 = st.columns([3,1])

with col1:
    st.subheader("1. Describe your presentation")
    with st.form("generate"):
        prompt = st.text_area(
            "Type your request:",
            height=120,
            placeholder="e.g., Create a 5-slide modern presentation about healthcare design."
        )
        phase = st.selectbox("Phase (optional):", ["Any","Plan","Design","Build","Test","Support"])
        submit = st.form_submit_button("Generate Presentation")

    if submit:
        tags = []
        if phase and phase != "Any":
            tags.append(phase)

        status = st.empty()
        progress = st.progress(0)
        try:
            status.text("Generating presentation...")
            progress.progress(20)
            out_path, log = generate_presentation(prompt=prompt, style=phase, tag_filters=tags)
            progress.progress(90)
            status.text("Done. Preparing preview...")
            with open(out_path, "rb") as f:
                st.download_button("📥 Download Generated PPTX", f, file_name=os.path.basename(out_path))
            progress.progress(100)
            status.success("Presentation generated successfully.")

            st.subheader("Slide previews")
            prs = Presentation(out_path)
            for i, slide in enumerate(prs.slides[:5]):
                title = ""
                try:
                    title = slide.shapes.title.text
                except Exception:
                    title = f"Slide {i+1}"
                img = Image.new("RGB", (720, 405), color=(245,245,245))
                d = ImageDraw.Draw(img)
                try:
                    font = ImageFont.load_default()
                except Exception:
                    font = None
                d.text((20,20), title, fill=(10,10,10), font=font)
                st.image(img, caption=f"Slide {i+1}: {title}")
        except Exception as e:
            logger.exception("Generation failed")
            st.error(f"Error: {e}")

with col2:
    st.subheader("Generated Presentations")
    generated_files = list_generated_presentations()
    if generated_files:
        for g in generated_files:
            st.write(f"📄 {g}")
    else:
        st.caption("No previously generated PPTs found.")
